"""
Document parser module for EduSum.
Handles multi-format document ingestion and text extraction.
"""

from abc import ABC, abstractmethod
from typing import Dict, List, Optional
import pymupdf  # PyMuPDF for PDF
from pptx import Presentation
from docx import Document


class DocumentParser(ABC):
    """Abstract base class for document parsers."""

    @abstractmethod
    def extract_text(self, file_path: str) -> str:
        """Extract text from document."""
        pass

    @abstractmethod
    def extract_metadata(self, file_path: str) -> Dict:
        """Extract metadata (structure, headings, etc.)."""
        pass


class PDFParser(DocumentParser):
    """Parser for PDF files using PyMuPDF."""

    def extract_text(self, file_path: str) -> str:
        """
        Extract text from PDF while preserving reading order.
        
        Args:
            file_path: Path to PDF file.
            
        Returns:
            Extracted text.
        """
        text = ""
        try:
            doc = pymupdf.open(file_path)
            for page_num in range(len(doc)):
                page = doc[page_num]
                # Use get_text() for better ordering in multi-column layouts
                text += page.get_text("text")
                text += "\n---PAGE BREAK---\n"
            doc.close()
        except Exception as e:
            raise ValueError(f"Error extracting PDF: {e}")
        return text

    def extract_metadata(self, file_path: str) -> Dict:
        """Extract PDF metadata and structure."""
        metadata = {
            "file_type": "pdf",
            "pages": 0,
            "title": "",
            "author": ""
        }
        try:
            doc = pymupdf.open(file_path)
            metadata["pages"] = len(doc)
            md = doc.metadata
            if md:
                metadata["title"] = md.get("title", "")
                metadata["author"] = md.get("author", "")
            doc.close()
        except Exception as e:
            print(f"Error extracting PDF metadata: {e}")
        return metadata


class PPTXParser(DocumentParser):
    """Parser for PowerPoint files."""

    def extract_text(self, file_path: str) -> str:
        """
        Extract text from PPTX slides and speaker notes.
        
        Args:
            file_path: Path to PPTX file.
            
        Returns:
            Extracted text.
        """
        text = ""
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n---SLIDE {slide_num + 1}---\n"
                
                # Extract text from slides
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                
                # Extract speaker notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes:
                        text += f"\n[NOTES: {notes}]\n"
        except Exception as e:
            raise ValueError(f"Error extracting PPTX: {e}")
        return text

    def extract_metadata(self, file_path: str) -> Dict:
        """Extract PPTX metadata."""
        metadata = {
            "file_type": "pptx",
            "slides": 0,
            "title": ""
        }
        try:
            prs = Presentation(file_path)
            metadata["slides"] = len(prs.slides)
            if prs.core_properties:
                metadata["title"] = prs.core_properties.title or ""
        except Exception as e:
            print(f"Error extracting PPTX metadata: {e}")
        return metadata


class DOCXParser(DocumentParser):
    """Parser for Word documents."""

    def extract_text(self, file_path: str) -> str:
        """
        Extract text from DOCX while preserving structure.
        
        Args:
            file_path: Path to DOCX file.
            
        Returns:
            Extracted text.
        """
        text = ""
        try:
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            # Also extract from tables
            for table in doc.tables:
                text += "\n---TABLE---\n"
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    text += row_text + "\n"
        except Exception as e:
            raise ValueError(f"Error extracting DOCX: {e}")
        return text

    def extract_metadata(self, file_path: str) -> Dict:
        """Extract DOCX metadata."""
        metadata = {
            "file_type": "docx",
            "paragraphs": 0,
            "title": ""
        }
        try:
            doc = Document(file_path)
            metadata["paragraphs"] = len(doc.paragraphs)
            if doc.core_properties:
                metadata["title"] = doc.core_properties.title or ""
        except Exception as e:
            print(f"Error extracting DOCX metadata: {e}")
        return metadata


class TXTParser(DocumentParser):
    """Parser for plain text files."""

    def extract_text(self, file_path: str) -> str:
        """
        Extract text from TXT file.
        
        Args:
            file_path: Path to TXT file.
            
        Returns:
            File text content.
        """
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            raise ValueError(f"Error reading TXT file: {e}")

    def extract_metadata(self, file_path: str) -> Dict:
        """Extract TXT metadata."""
        import os
        metadata = {
            "file_type": "txt",
            "title": os.path.basename(file_path),
            "size_bytes": 0
        }
        try:
            metadata["size_bytes"] = os.path.getsize(file_path)
        except Exception as e:
            print(f"Error extracting TXT metadata: {e}")
        return metadata


def detect_file_type(file_path: str) -> str:
    """
    Detect file type using multiple strategies (not just extension).
    PRD: "System must detect file type automatically without relying on file extension."
    
    Falls back gracefully: mimetypes → extension.
    
    Args:
        file_path: Path to file.
        
    Returns:
        Detected file type string (pdf, pptx, docx, txt).
    """
    import mimetypes

    # Strategy 1: Use mimetypes library (built-in, cross-platform)
    mime_type, _ = mimetypes.guess_type(file_path)
    mime_to_type = {
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "text/plain": "txt",
    }
    if mime_type and mime_type in mime_to_type:
        return mime_to_type[mime_type]

    # Strategy 2: Check magic bytes (file signature)
    try:
        with open(file_path, 'rb') as f:
            header = f.read(8)
        # PDF starts with %PDF
        if header[:4] == b'%PDF':
            return "pdf"
        # PPTX/DOCX are ZIP files (PK signature), differentiate by extension
        if header[:2] == b'PK':
            ext = file_path.split(".")[-1].lower()
            if ext in ("pptx", "docx"):
                return ext
            # Try to peek inside the ZIP to determine type
            import zipfile
            try:
                with zipfile.ZipFile(file_path, 'r') as zf:
                    names = zf.namelist()
                    if any('ppt/' in n or 'ppt\\' in n for n in names):
                        return "pptx"
                    if any('word/' in n or 'word\\' in n for n in names):
                        return "docx"
            except Exception:
                pass
    except Exception:
        pass

    # Strategy 3: Fall back to extension
    ext = file_path.split(".")[-1].lower()
    if ext in ("pdf", "pptx", "docx", "txt"):
        return ext

    return ext


class DocumentParserFactory:
    """Factory for selecting the right parser based on file type."""

    _parsers = {
        "pdf": PDFParser,
        "pptx": PPTXParser,
        "docx": DOCXParser,
        "txt": TXTParser
    }

    @staticmethod
    def get_parser(file_path: str) -> DocumentParser:
        """
        Get appropriate parser for file type.
        Uses auto-detection (PRD requirement) with extension fallback.
        
        Args:
            file_path: Path to file.
            
        Returns:
            Parser instance.
        """
        file_type = detect_file_type(file_path)
        parser_class = DocumentParserFactory._parsers.get(file_type)
        
        if not parser_class:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        return parser_class()


def parse_document(file_path: str) -> Dict:
    """
    Convenience function to parse any supported document.
    
    Args:
        file_path: Path to document.
        
    Returns:
        Dictionary with extracted text and metadata.
    """
    parser = DocumentParserFactory.get_parser(file_path)
    return {
        "text": parser.extract_text(file_path),
        "metadata": parser.extract_metadata(file_path)
    }
